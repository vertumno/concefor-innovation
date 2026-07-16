from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT = Path(__file__).resolve().parent / "opcoes"
OUT.mkdir(parents=True, exist_ok=True)

LOGO_CONCEFOR = ROOT / "contexto/_marca/logos/concefor/Arquivos em PNG - alta resolução/VIII_Concefor_logo_horizontal_branco.png"
ASSINATURA_BRANCA = ROOT / "contexto/_marca/logos/20-anos-cefor/PNG/Selo_20_anos_Cefor_Aplicacao_Cefor_horizontal_horizontal_branco.png"
ASSINATURA_COLORIDA = ROOT / "contexto/_marca/logos/20-anos-cefor/PNG/Selo_20_anos_Cefor_Aplicacao_Cefor_horizontal_horizontal_colorida.png"
FONT_DIR = Path(__file__).resolve().parent / "assets/fontes"
TEKO_TTF = FONT_DIR / "Teko-wght.ttf"
MONTSERRAT_TTF = FONT_DIR / "Montserrat-wght.ttf"

NAVY = RGBColor(12, 42, 105)
BLUE = RGBColor(31, 74, 161)
BLUE_2 = RGBColor(20, 91, 154)
CYAN = RGBColor(0, 220, 201)
CYAN_DARK = RGBColor(0, 176, 170)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(13, 44, 93)
MUTED = RGBColor(196, 222, 237)


def box(slide, x, y, w, h, color, radius=False, line=None, line_width=1):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def line(slide, x1, y1, x2, y2, color=CYAN, width=1.5):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x1), Inches(y1), Inches(x2-x1), Inches(y2-y1))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def text(slide, value, x, y, w, h, size, color=WHITE, font="Arial", bold=False,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0, tracking=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = 0.92 if font == "Bahnschrift SemiCondensed" else 1.02
    run = p.add_run()
    run.text = value
    # Tipografia oficial do design system: Teko nos títulos, Montserrat no apoio.
    run.font.name = "Teko" if font == "Bahnschrift SemiCondensed" else "Montserrat"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if tracking is not None:
        run._r.get_or_add_rPr().set('spc', str(tracking))
    return tb


def add_logo(slide, path, x, y, w):
    # Inserção integral do PNG oficial, com proporção preservada e sem recorte.
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


def circle_number(slide, n, x, y, diameter=0.34, fill=CYAN, color=INK):
    c = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background()
    text(slide, str(n), x, y+0.005, diameter, diameter-0.01, 12, color, "Arial", True,
         PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


prs = Presentation()
prs.slide_width = Inches(8)
prs.slide_height = Inches(10)
blank = prs.slide_layouts[6]

# OPÇÃO A — urgência editorial, data protagonista
s = prs.slides.add_slide(blank)
box(s, 0, 0, 8, 10, NAVY)
box(s, 0, 0, 8, 0.12, CYAN)
box(s, 0, 0.12, 2.45, 9.88, BLUE_2)
box(s, 0.23, 6.75, 2.05, 2.55, BLUE)
line(s, 2.45, 0.65, 2.48, 9.25, CYAN, 1)

add_logo(s, LOGO_CONCEFOR, 0.48, 0.5, 2.55)
text(s, "COMUNICADO AOS AUTORES", 4.68, 0.58, 2.78, 0.34, 11, CYAN, "Arial", True, PP_ALIGN.RIGHT)
text(s, "VERSÃO FINAL", 2.92, 1.55, 4.44, 0.55, 20, CYAN, "Arial", True)
text(s, "FALTA", 2.88, 2.0, 4.5, 0.72, 38, WHITE, "Bahnschrift SemiCondensed", True)
text(s, "UM PASSO.", 2.88, 2.66, 4.5, 0.72, 38, WHITE, "Bahnschrift SemiCondensed", True)
text(s, "Envie a versão final do seu trabalho aprovado pela Even3.", 2.94, 3.82, 4.15, 0.85, 18, WHITE, "Arial", False)

text(s, "PRAZO", 0.48, 2.05, 1.5, 0.3, 11, MUTED, "Arial", True)
text(s, "03", 0.42, 2.30, 1.7, 1.25, 56, WHITE, "Bahnschrift SemiCondensed", True)
line(s, 0.48, 3.55, 2.08, 3.58, CYAN, 1.5)
text(s, "AGO", 0.48, 3.68, 1.45, 0.55, 25, CYAN, "Bahnschrift SemiCondensed", True)
text(s, "2026", 0.48, 4.15, 1.45, 0.45, 16, WHITE, "Arial", True)

box(s, 2.92, 5.0, 4.42, 1.9, BLUE, True, CYAN, 1.2)
circle_number(s, 1, 3.18, 5.28)
text(s, "Use o template oficial do VIII Concefor", 3.66, 5.25, 3.35, 0.4, 14, WHITE, "Arial", True)
circle_number(s, 2, 3.18, 5.87)
text(s, "Faça o envio pela plataforma Even3", 3.66, 5.84, 3.35, 0.4, 14, WHITE, "Arial", True)
circle_number(s, 3, 3.18, 6.46)
text(s, "Confirme a inscrição de pelo menos um autor", 3.66, 6.38, 3.35, 0.55, 13, WHITE, "Arial", True)

text(s, "Para publicação nos Anais, todos os autores também devem constar no cadastro do trabalho.", 2.94, 7.22, 4.15, 0.72, 12, MUTED, "Arial")
line(s, 2.94, 8.05, 7.35, 8.08, CYAN, 1)
text(s, "17—20 AGO 2026  ·  VITÓRIA—ES", 2.94, 8.28, 4.15, 0.35, 11, WHITE, "Arial", True)
text(s, "concefor.cefor.ifes.edu.br", 2.94, 8.68, 4.15, 0.35, 14, CYAN, "Arial", True)
add_logo(s, ASSINATURA_BRANCA, 0.45, 8.98, 7.05)

# OPÇÃO B — narrativa institucional e hierarquia em módulos
s = prs.slides.add_slide(blank)
box(s, 0, 0, 8, 10, BLUE)
box(s, 0, 0, 8, 0.12, CYAN)
box(s, 5.95, 0.12, 2.05, 8.25, NAVY)
box(s, 0, 8.37, 8, 1.63, WHITE)
box(s, 0.52, 0.62, 0.09, 1.05, CYAN)
add_logo(s, LOGO_CONCEFOR, 0.78, 0.5, 2.55)

text(s, "TRABALHO", 0.52, 2.0, 4.95, 0.68, 35, WHITE, "Bahnschrift SemiCondensed", True)
text(s, "APROVADO?", 0.52, 2.64, 4.95, 0.68, 35, WHITE, "Bahnschrift SemiCondensed", True)
text(s, "GARANTA A PUBLICAÇÃO NOS ANAIS.", 0.55, 3.62, 4.95, 0.72, 18, CYAN, "Arial", True)
text(s, "O envio final deve ser feito pela Even3, usando o template oficial do evento.", 0.55, 4.45, 4.72, 0.88, 16, WHITE, "Arial")

for idx, (title_, body_) in enumerate([
    ("TEMPLATE", "formate a versão final"),
    ("EVEN3", "envie o arquivo na plataforma"),
    ("INSCRIÇÃO", "ao menos um autor inscrito"),
]):
    yy = 5.62 + idx * 0.82
    circle_number(s, idx+1, 0.56, yy, 0.38, CYAN, INK)
    text(s, title_, 1.10, yy-0.02, 1.35, 0.30, 12, CYAN, "Arial", True)
    text(s, body_, 2.43, yy-0.02, 3.05, 0.40, 12, WHITE, "Arial")
    if idx < 2:
        line(s, 1.10, yy+0.50, 5.45, yy+0.515, RGBColor(71, 126, 187), 0.7)

text(s, "PRAZO FINAL", 6.25, 1.18, 1.45, 0.35, 11, CYAN, "Arial", True, PP_ALIGN.CENTER)
text(s, "03", 6.05, 1.78, 1.84, 1.12, 54, WHITE, "Bahnschrift SemiCondensed", True, PP_ALIGN.CENTER)
line(s, 6.28, 2.94, 7.66, 2.98, CYAN, 2)
text(s, "08", 6.05, 3.15, 1.84, 1.08, 46, WHITE, "Bahnschrift SemiCondensed", True, PP_ALIGN.CENTER)
text(s, "2026", 6.12, 4.28, 1.70, 0.38, 15, MUTED, "Arial", True, PP_ALIGN.CENTER)
box(s, 6.25, 5.13, 1.48, 0.58, CYAN, True)
text(s, "PELA EVEN3", 6.25, 5.13, 1.48, 0.58, 11, INK, "Arial", True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
text(s, "17—20\nAGOSTO\nVITÓRIA—ES", 6.15, 6.55, 1.68, 1.20, 15, WHITE, "Bahnschrift SemiCondensed", True, PP_ALIGN.CENTER)

text(s, "concefor.cefor.ifes.edu.br", 0.55, 8.65, 3.45, 0.38, 15, INK, "Arial", True)
text(s, "Orientações e templates no site oficial", 0.55, 9.10, 3.45, 0.35, 11, BLUE, "Arial")
add_logo(s, ASSINATURA_COLORIDA, 4.25, 8.63, 3.20)

dest = OUT / "2026-07-16_pre_card_prazo-envio-final_opcoes-editaveis.pptx"
prs.save(dest)
print(dest)


# Renderização PNG determinística (1080×1350), sem regenerar ou redesenhar logos.
from PIL import Image, ImageDraw, ImageFont

W, H = 1080, 1350
S = 135.0  # pixels por polegada no layout 8×10
P_NAVY = (12, 42, 105)
P_BLUE = (31, 74, 161)
P_BLUE_2 = (20, 91, 154)
P_CYAN = (0, 220, 201)
P_WHITE = (255, 255, 255)
P_INK = (13, 44, 93)
P_MUTED = (196, 222, 237)
P_LINE = (71, 126, 187)


def px(v): return int(round(v * S))


def font(size, bold=False, condensed=False):
    path = TEKO_TTF if condensed else MONTSERRAT_TTF
    fnt = ImageFont.truetype(str(path), size=int(size * 1.88))
    fnt.set_variation_by_axes([700 if bold else 400])
    return fnt


def fit_logo(canvas, path, box_):
    # O arquivo inteiro é escalado proporcionalmente para caber; não há crop nem recoloração.
    logo = Image.open(path).convert("RGBA")
    x, y, w, h = box_
    ratio = min(w / logo.width, h / logo.height)
    logo = logo.resize((round(logo.width * ratio), round(logo.height * ratio)), Image.Resampling.LANCZOS)
    xx = x + (w - logo.width) // 2
    yy = y + (h - logo.height) // 2
    canvas.alpha_composite(logo, (xx, yy))


def draw_text(draw, xy, value, fnt, fill, anchor="la", spacing=4, align="left"):
    draw.multiline_text(xy, value, font=fnt, fill=fill, anchor=anchor, spacing=spacing, align=align)


def rounded(draw, box_, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(box_, radius=radius, fill=fill, outline=outline, width=width)


def render_a():
    im = Image.new("RGBA", (W, H), P_NAVY + (255,)); d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 16), fill=P_CYAN)
    d.rectangle((0, 16, px(2.45), H), fill=P_BLUE_2)
    d.polygon([(0, 790), (px(2.45), 690), (px(2.45), 1120), (0, 1240)], fill=P_BLUE)
    d.rectangle((px(2.45), 88, px(2.48), 1250), fill=P_CYAN)
    fit_logo(im, LOGO_CONCEFOR, (px(.40), px(.40), px(2.35), px(1.06)))
    draw_text(d, (1005, 84), "COMUNICADO AOS AUTORES", font(10, True), P_CYAN, "ra")

    draw_text(d, (px(2.92), px(1.55)), "VERSÃO FINAL", font(16, True), P_CYAN)
    draw_text(d, (px(2.90), px(2.02)), "FALTA", font(34, True, True), P_WHITE)
    draw_text(d, (px(2.90), px(2.64)), "UM PASSO.", font(34, True, True), P_WHITE)
    draw_text(d, (px(2.94), px(3.88)), "Envie a versão final do seu\ntrabalho aprovado pela Even3.", font(14), P_WHITE, spacing=8)

    draw_text(d, (px(.48), px(2.05)), "PRAZO", font(10, True), P_MUTED)
    draw_text(d, (px(.42), px(2.28)), "03", font(53, True, True), P_WHITE)
    d.rectangle((px(.48), px(3.55), px(2.08), px(3.58)), fill=P_CYAN)
    draw_text(d, (px(.48), px(3.68)), "AGO", font(24, True, True), P_CYAN)
    draw_text(d, (px(.48), px(4.17)), "2026", font(14, True), P_WHITE)

    panel = (px(2.92), px(5.0), px(7.34), px(6.98))
    rounded(d, panel, 22, P_BLUE, P_CYAN, 2)
    items = [
        "Use o template oficial do VIII Concefor",
        "Faça o envio pela plataforma Even3",
        "Confirme a inscrição de\npelo menos um autor",
    ]
    for i, label in enumerate(items):
        cy = px(5.31 + i*.61)
        d.ellipse((px(3.17), cy-22, px(3.17)+44, cy+22), fill=P_CYAN)
        draw_text(d, (px(3.17)+22, cy), str(i+1), font(10, True), P_INK, "mm")
        draw_text(d, (px(3.68), cy), label, font(10.8, True), P_WHITE, "lm", spacing=2)

    draw_text(d, (px(2.94), px(7.25)), "Para publicação nos Anais, todos os autores também\ndevem constar no cadastro do trabalho.", font(10.5), P_MUTED, spacing=5)
    d.rectangle((px(2.94), px(8.05), px(7.35), px(8.07)), fill=P_CYAN)
    draw_text(d, (px(2.94), px(8.30)), "17—20 AGO 2026  ·  VITÓRIA—ES", font(10.5, True), P_WHITE)
    draw_text(d, (px(2.94), px(8.67)), "concefor.cefor.ifes.edu.br", font(13, True), P_CYAN)
    fit_logo(im, ASSINATURA_BRANCA, (px(.42), px(9.00), px(7.10), px(.76)))
    path = OUT / "2026-07-16_pre_card_prazo-envio-final_opcao-a.png"
    im.convert("RGB").save(path, quality=96)
    return path


def render_b():
    im = Image.new("RGBA", (W, H), P_BLUE + (255,)); d = ImageDraw.Draw(im)
    d.rectangle((0, 0, W, 16), fill=P_CYAN)
    d.rectangle((px(5.95), 16, W, px(8.37)), fill=P_NAVY)
    d.rectangle((0, px(8.37), W, H), fill=P_WHITE)
    d.polygon([(0, 780), (650, 620), (805, 1110), (0, 1170)], fill=(25, 83, 158))
    d.rectangle((px(.52), px(.62), px(.61), px(1.67)), fill=P_CYAN)
    fit_logo(im, LOGO_CONCEFOR, (px(.70), px(.38), px(2.70), px(1.22)))

    draw_text(d, (px(.52), px(2.02)), "TRABALHO", font(32, True, True), P_WHITE)
    draw_text(d, (px(.52), px(2.60)), "APROVADO?", font(32, True, True), P_WHITE)
    draw_text(d, (px(.55), px(3.66)), "GARANTA A PUBLICAÇÃO NOS ANAIS.", font(14, True), P_CYAN)
    draw_text(d, (px(.55), px(4.46)), "O envio final deve ser feito pela Even3,\nusando o template oficial do evento.", font(13), P_WHITE, spacing=7)

    data = [
        ("TEMPLATE", "formate a versão final"),
        ("EVEN3", "envie o arquivo na plataforma"),
        ("INSCRIÇÃO", "ao menos um autor inscrito"),
    ]
    for i, (ttl, body) in enumerate(data):
        cy = px(5.70 + i*.78)
        d.ellipse((px(.55), cy-24, px(.55)+48, cy+24), fill=P_CYAN)
        draw_text(d, (px(.55)+24, cy), str(i+1), font(10, True), P_INK, "mm")
        draw_text(d, (px(1.05), cy), ttl, font(11, True), P_CYAN, "lm")
        draw_text(d, (px(2.30), cy), body, font(10.5), P_WHITE, "lm")
        if i < 2: d.rectangle((px(1.05), cy+48, px(5.43), cy+50), fill=P_LINE)

    draw_text(d, (px(6.97), px(1.22)), "PRAZO FINAL", font(10, True), P_CYAN, "ma")
    draw_text(d, (px(6.97), px(1.74)), "03", font(51, True, True), P_WHITE, "ma")
    d.rectangle((px(6.28), px(2.95), px(7.66), px(2.98)), fill=P_CYAN)
    draw_text(d, (px(6.97), px(3.16)), "08", font(44, True, True), P_WHITE, "ma")
    draw_text(d, (px(6.97), px(4.26)), "2026", font(14, True), P_MUTED, "ma")
    rounded(d, (px(6.23), px(5.10), px(7.72), px(5.70)), 16, P_CYAN)
    draw_text(d, (px(6.975), px(5.40)), "PELA EVEN3", font(10.5, True), P_INK, "mm")
    draw_text(d, (px(6.97), px(6.55)), "17—20\nAGOSTO\nVITÓRIA—ES", font(14, True, True), P_WHITE, "ma", spacing=5, align="center")

    draw_text(d, (px(.55), px(8.65)), "concefor.cefor.ifes.edu.br", font(13.5, True), P_INK)
    draw_text(d, (px(.55), px(9.10)), "Orientações e templates no site oficial", font(10.5), P_BLUE)
    fit_logo(im, ASSINATURA_COLORIDA, (px(4.18), px(8.55), px(3.40), px(1.05)))
    path = OUT / "2026-07-16_pre_card_prazo-envio-final_opcao-b.png"
    im.convert("RGB").save(path, quality=96)
    return path


print(render_a())
print(render_b())
